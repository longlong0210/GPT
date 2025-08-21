import os
import google.generativeai as genai
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation

# ==========================
# 1. Cấu hình Google AI (Giữ nguyên)
# ==========================
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    # Vui lòng thay thế bằng API key của bạn nếu không dùng biến môi trường
    GOOGLE_API_KEY = "AIzaSyABLonsDEQ7veJFWZf6lLlHvtPw9K4lBMs"

genai.configure(api_key=GOOGLE_API_KEY)
print("✅ Đã cấu hình Google AI thành công.")


# ==========================
# 2. Hàm trích xuất nội dung (Giữ nguyên)
# ==========================
def extract_text_from_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text_from_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    texts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows():
            row_data = [str(cell.value) for cell in row if cell.value]
            if row_data:
                texts.append(" ".join(row_data))
    return "\n".join(texts)

def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

extractors = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".pptx": extract_text_from_pptx,
    ".xlsx": extract_text_from_xlsx,
    ".txt": extract_text_from_txt
}

def read_file_content(file_path):
    """Tự phát hiện loại file và đọc nội dung"""
    if not os.path.exists(file_path):
        return f"Lỗi: File '{file_path}' không tồn tại."
    
    ext = os.path.splitext(file_path)[1].lower()
    if ext in extractors:
        try:
            return extractors[ext](file_path)
        except Exception as e:
            return f"Lỗi khi đọc file {file_path}: {e}"
    else:
        return f"Lỗi: Định dạng file '{ext}' chưa được hỗ trợ."

# ==========================
# 3. Gọi Google AI (Giữ nguyên)
# ==========================
def analyze_files_with_gemini(file_paths, user_prompt):
    """Đọc nhiều file và gửi nội dung cho Google AI"""
    combined_text = ""
    print("\n⏳ Đang đọc và trích xuất nội dung từ các file...")
    for file_path in file_paths:
        content = read_file_content(file_path)
        if content.startswith("Lỗi:"):
            print(f"   - {content}")
            continue
        print(f"   - Đã đọc thành công file: {os.path.basename(file_path)}")
        combined_text += f"\n\n--- Nội dung từ file {os.path.basename(file_path)} ---\n{content}"

    if not combined_text.strip():
        return "Không có dữ liệu hợp lệ để phân tích."

    model = genai.GenerativeModel("gemini-1.5-flash") # Sử dụng model mới hơn
    prompt = f"{user_prompt}\n\nDữ liệu từ các file:\n{combined_text}"

    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Lỗi khi gọi Google AI: {e}"

# =========================================================
# 4. Chạy chương trình với đầu vào từ người dùng
# =========================================================
if __name__ == "__main__":
    
    # --- PHẦN 1: NHẬP ĐƯỜNG DẪN FILE ---
    file_paths = []
    print("\nNhập đường dẫn đến các file bạn muốn phân tích.")
    print("Gõ 'xong' hoặc 'done' trên một dòng riêng để kết thúc.")
    
    while True:
        path = input(f"Đường dẫn file {len(file_paths) + 1}: ").strip()
        
        # Kiểm tra lệnh dừng
        if path.lower() in ['xong', 'done']:
            break
        
        # Kiểm tra xem file có tồn tại không trước khi thêm vào danh sách
        if os.path.exists(path):
            file_paths.append(path)
            print(f"-> Đã thêm file: {os.path.basename(path)}")
        else:
            print("-> Lỗi: File không tồn tại. Vui lòng kiểm tra lại đường dẫn.")

    # --- PHẦN 2: NHẬP CÂU HỎI VÀ GỌI AI ---
    if not file_paths:
        print("\nBạn chưa chọn file nào. Chương trình kết thúc.")
    else:
        user_prompt = input("\nNhập câu hỏi hoặc yêu cầu của bạn về các tài liệu trên:\n")
        
        print("\n🔍 Đang phân tích file...")
        result = analyze_files_with_gemini(file_paths, user_prompt)
        print("\n" + "="*50)
        print("✅ Kết quả từ Google AI:")
        print("="*50)
        print(result)