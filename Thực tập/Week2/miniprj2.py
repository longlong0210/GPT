import os
import re
import unicodedata
import google.generativeai as genai
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
import pytesseract
from PIL import Image
import io

# =========================================================
# PHẦN 1: CẤU HÌNH VÀ CÁC HẰNG SỐ
# =========================================================

# --- Cấu hình Google AI ---
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    # VUI LÒNG THAY THẾ BẰNG API KEY CỦA BẠN
    GOOGLE_API_KEY = "AIzaSyABLonsDEQ7veJFWZf6lLlHvtPw9K4lBMs"

try:
    genai.configure(api_key=GOOGLE_API_KEY)
    print("✅ Đã cấu hình Google AI thành công.")
except Exception as e:
    print(f"❌ Lỗi cấu hình Google AI: {e}")

# --- Danh sách từ dừng Tiếng Việt ---
VIETNAMESE_STOP_WORDS = [
    "và", "là", "mà", "thì", "của", "ở", "tại", "bị", "bởi", "cả", "các", "cái", "cần",
    "cho", "chứ", "chưa", "có", "cũng", "đã", "đang", "đây", "để", "đến", "đều", "điều",
    "do", "đó", "được", "khi", "không", "lại", "lên", "lúc", "mỗi", "một", "nên", "nếu",
    "ngay", "như", "nhưng", "những", "nơi", "nữa", "phải", "qua", "ra", "rằng", "rất",
    "rồi", "sau", "sẽ", "so", "sự", "tại", "theo", "thì", "trên", "trước", "từ", "vẫn",
    "vào", "vậy", "về", "vì", "việc", "với", "vừa"
]

# =========================================================
# PHẦN 2: CÁC HÀM XỬ LÝ DỮ LIỆU
# =========================================================

# --- Các hàm trích xuất nội dung từ file ---

def extract_text_from_pdf(file_path):
    """Trích xuất văn bản từ PDF, tự động dùng OCR nếu cần."""
    text = ""
    try:
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        
        if len(text.strip()) < 100:
            print(f"   - Văn bản trong file '{os.path.basename(file_path)}' ít. Thử dùng OCR...")
            ocr_text = ""
            with fitz.open(file_path) as doc:
                for page_num, page in enumerate(doc):
                    pix = page.get_pixmap(dpi=300)
                    img_bytes = pix.tobytes("png")
                    image = Image.open(io.BytesIO(img_bytes))
                    page_text = pytesseract.image_to_string(image, lang='vie+eng')
                    ocr_text += page_text + "\n"
            return ocr_text
    except Exception as e:
        return f"Lỗi khi xử lý PDF '{file_path}': {e}"
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def extract_text_from_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    texts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            row_data = [str(cell.value) for cell in row if cell.value]
            if row_data:
                texts.append(" ".join(row_data))
    return "\n".join(texts)

def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

# --- Hàm tổng hợp: Đọc file ---

def read_file_content(file_path):
    """Tự phát hiện loại file và đọc nội dung."""
    extractors = {
        ".pdf": extract_text_from_pdf, ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx, ".xlsx": extract_text_from_xlsx,
        ".txt": extract_text_from_txt
    }
    ext = os.path.splitext(file_path)[1].lower()
    if ext in extractors:
        try:
            return extractors[ext](file_path)
        except Exception as e:
            return f"Lỗi khi đọc file {os.path.basename(file_path)}: {e}"
    return f"Lỗi: Định dạng file '{ext}' chưa được hỗ trợ."

# --- Các hàm làm sạch và phân đoạn ---

def normalize_and_clean_text(text: str) -> str:
    """Làm sạch và chuẩn hóa một đoạn văn bản."""
    text = unicodedata.normalize('NFC', text)
    text = text.lower()
    text = re.sub(r'https?://\S+|www\.\S+', '', text)
    text = re.sub(r'<.*?>', '', text)
    text = re.sub(r'[^\w\s]', '', text)
    words = text.split()
    filtered_words = [word for word in words if word not in VIETNAMESE_STOP_WORDS]
    text = " ".join(filtered_words)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def chunk_text_by_sentence(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[str]:
    """Phân đoạn văn bản một cách thông minh, giữ trọn vẹn câu."""
    sentences = re.split(r'(?<=[.?!])\s+', text)
    sentences = [s.strip() for s in sentences if s.strip()]
    if not sentences: return []

    chunks = []
    current_chunk = ""
    for sentence in sentences:
        if len(current_chunk) + len(sentence) + 1 <= chunk_size:
            current_chunk += sentence + " "
        else:
            if current_chunk: chunks.append(current_chunk.strip())
            current_chunk = sentence + " "
    if current_chunk: chunks.append(current_chunk.strip())
    return chunks

# =========================================================
# PHẦN 3: HÀM TỔNG HỢP GỌI AI
# =========================================================
def process_and_analyze(file_paths, user_prompt):
    """
    Thực hiện toàn bộ quy trình: Đọc -> Làm sạch -> Phân đoạn -> Gọi AI.
    """
    all_chunks = []
    print("\n⏳ Bắt đầu quy trình xử lý file...")
    for file_path in file_paths:
        # 1. Đọc file
        print(f"   - [1/3] Đang đọc file: {os.path.basename(file_path)}")
        raw_content = read_file_content(file_path)
        if raw_content.startswith("Lỗi:"):
            print(f"     -> {raw_content}")
            continue
        
        # 2. Làm sạch
        print("   - [2/3] Đang làm sạch và chuẩn hóa nội dung...")
        cleaned_content = normalize_and_clean_text(raw_content)
        
        # 3. Phân đoạn
        print("   - [3/3] Đang phân đoạn văn bản...")
        chunks = chunk_text_by_sentence(cleaned_content)
        all_chunks.extend(chunks)
        print(f"     -> Hoàn tất, tạo ra {len(chunks)} đoạn.")

    if not all_chunks:
        return "Không có dữ liệu hợp lệ để phân tích sau khi xử lý."

    combined_text = "\n\n---\n\n".join(all_chunks)
    
    print("\n🧠 Đang gửi dữ liệu đã xử lý đến Google AI...")
    model = genai.GenerativeModel("gemini-2.5-flash")
    prompt = f"Dựa trên các đoạn văn bản dưới đây, hãy thực hiện yêu cầu sau: '{user_prompt}'.\n\nDữ liệu:\n{combined_text}"

    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Lỗi khi gọi Google AI: {e}"

# =========================================================
# PHẦN 4: CHẠY CHƯƠNG TRÌNH CHÍNH
# =========================================================
if __name__ == "__main__":
    file_paths = []
    print("\nNhập đường dẫn đến các file bạn muốn phân tích.")
    print("Gõ 'xong' hoặc 'done' trên một dòng riêng để kết thúc.")
    
    while True:
        path = input(f"Đường dẫn file {len(file_paths) + 1}: ").strip()
        if path.lower() in ['xong', 'done']: break
        if os.path.exists(path):
            file_paths.append(path)
            print(f"-> Đã thêm file: {os.path.basename(path)}")
        else:
            print("-> Lỗi: File không tồn tại. Vui lòng kiểm tra lại đường dẫn.")

    if file_paths:
        user_prompt = input("\nNhập yêu cầu của bạn về các tài liệu trên (ví dụ: tóm tắt nội dung):\n")
        
        result = process_and_analyze(file_paths, user_prompt)
        
        print("\n" + "="*50)
        print("✅ KẾT QUẢ PHÂN TÍCH TỪ GOOGLE AI:")
        print("="*50)
        print(result)
    else:
        print("\nBạn chưa chọn file nào. Chương trình kết thúc.")