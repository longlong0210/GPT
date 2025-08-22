import os
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation

# ==== HÀM ĐỌC FILE NHIỀU ĐỊNH DẠNG ====
def read_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    content = []

    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            content = [line.strip() for line in f if line.strip()]

    elif ext == ".docx":
        doc = Document(file_path)
        content = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    elif ext == ".pdf":
        pdf = fitz.open(file_path)
        for page in pdf:
            text = page.get_text().strip()
            if text:
                content.extend(text.split("\n"))

    elif ext == ".pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text:
                        content.append(text)

    else:
        raise ValueError(f"Không hỗ trợ định dạng {ext}")

    return content

# ==== CHƯƠNG TRÌNH CHÍNH ====
print("=== VECTOR SEARCH FAISS + VIETNAMESE EMBEDDING ===")

choice = input("Bạn muốn nhập dữ liệu thủ công (1) hay upload file (2)? ").strip()

if choice == "1":
    documents = []
    print("Nhập từng câu/đoạn (gõ 'xong' để kết thúc):")
    while True:
        line = input("> ").strip()
        if line.lower() == "xong":
            break
        if line:
            documents.append(line)
elif choice == "2":
    file_path = input("Nhập đường dẫn file: ").strip()
    documents = read_file(file_path)
else:
    raise ValueError("Lựa chọn không hợp lệ!")

if not documents:
    raise ValueError("Không có dữ liệu để xử lý!")

# 1. Tải mô hình
print("Đang tải mô hình embedding...")
model = SentenceTransformer('bkai-foundation-models/vietnamese-bi-encoder')
print("Tải mô hình thành công!")

# 2. Mã hóa vector
print("\nĐang tạo vector từ dữ liệu...")
document_embeddings = model.encode(documents)
print(f"Đã tạo {document_embeddings.shape[0]} vector, mỗi vector có {document_embeddings.shape[1]} chiều.")

# 3. Xây FAISS index
d = document_embeddings.shape[1]
index = faiss.IndexFlatL2(d)
index.add(document_embeddings.astype('float32'))
print(f"Đã thêm {index.ntotal} vector vào FAISS index.")

# 4. Người dùng hỏi
query = input("\nNhập câu hỏi của bạn: ").strip()
k = int(input("Số kết quả muốn tìm: "))

# 5. Encode câu hỏi
query_embedding = model.encode([query])

# 6. Tìm kiếm
distances, indices = index.search(query_embedding.astype('float32'), k)

# 7. Hiển thị kết quả
print("\n=== KẾT QUẢ TÌM KIẾM ===")
for i, idx in enumerate(indices[0]):
    print(f"\n--- Kết quả {i+1} ---")
    print(f"Nội dung: {documents[idx]}")
    print(f"Khoảng cách: {distances[0][i]:.4f}")
